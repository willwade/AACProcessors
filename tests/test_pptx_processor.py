import os
import tempfile
from collections.abc import Generator

import pytest
from pptx import Presentation  # type: ignore

from aac_processors.pptx_processor import PowerPointProcessor
from aac_processors.tree_structure import ButtonType


@pytest.fixture
def test_pptx_file() -> Generator[str, None, None]:
    """Create a test PowerPoint file."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        # Create test presentation
        prs = Presentation()

        # Add a slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add some text boxes
        left = top = width = height = prs.slide_width // 4

        # Add first text box
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.text = "Hello World"

        # Add second text box
        text_box2 = slide.shapes.add_textbox(left * 2, top * 2, width, height)
        tf = text_box2.text_frame
        tf.text = "Test Text"

        # Save the file
        prs.save(f.name)
        yield f.name

    # Cleanup
    if os.path.exists(f.name):
        os.remove(f.name)


def test_can_process() -> None:
    """Test file type detection."""
    processor = PowerPointProcessor()
    assert processor.can_process("test.pptx")
    assert processor.can_process("TEST.PPTX")
    assert not processor.can_process("test.doc")
    assert not processor.can_process("test.txt")


def test_load_into_tree(test_pptx_file: str) -> None:
    """Test loading PowerPoint into tree structure."""
    processor = PowerPointProcessor()
    tree = processor.load_into_tree(test_pptx_file)

    # Check basic tree structure
    assert len(tree.pages) == 1
    assert tree.root_id == "slide_1"

    # Check first page
    page = tree.pages["slide_1"]
    assert page.name == "Slide 1"
    assert len(page.buttons) == 2

    # Check buttons
    button1 = page.buttons[0]
    assert button1.label == "Hello World"
    assert button1.type == ButtonType.SPEAK
    assert button1.vocalization == "Hello World"

    button2 = page.buttons[1]
    assert button2.label == "Test Text"
    assert button2.type == ButtonType.SPEAK
    assert button2.vocalization == "Test Text"


def test_save_from_tree(test_pptx_file: str, tmp_path: str) -> None:
    """Test saving tree back to PowerPoint."""
    processor = PowerPointProcessor()
    tree = processor.load_into_tree(test_pptx_file)

    # Save to new file
    output_path = os.path.join(tmp_path, "output.pptx")
    processor.save_from_tree(tree, str(output_path))

    # Verify file was created
    assert os.path.exists(output_path)

    # Load saved file and verify content
    prs = Presentation(str(output_path))
    assert len(prs.slides) == 1

    slide = prs.slides[0]
    text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    assert len(text_shapes) == 2

    texts = [shape.text_frame.text for shape in text_shapes]
    assert "Hello World" in texts
    assert "Test Text" in texts


def test_process_files_extract(test_pptx_file: str) -> None:
    """Test extracting texts from PowerPoint."""
    processor = PowerPointProcessor()
    processor.set_source_file(test_pptx_file)

    # Extract texts
    processor.process_files("", None)

    # Verify extracted texts
    assert len(processor.collected_texts) == 2
    assert "Hello World" in processor.collected_texts
    assert "Test Text" in processor.collected_texts


def test_process_files_translate(test_pptx_file: str, tmp_path: str) -> None:
    """Test translating PowerPoint file."""
    processor = PowerPointProcessor()
    processor.set_source_file(test_pptx_file)

    # Create translations
    translations = {
        "Hello World": "Hola Mundo",
        "Test Text": "Texto de Prueba",
        "target_lang": "es",
    }

    # Process translation
    output_path = processor.process_files(str(tmp_path), translations)

    # Verify output file
    assert output_path is not None
    assert os.path.exists(output_path)
    assert output_path.endswith("_es.pptx")

    # Load translated file and verify content
    prs = Presentation(output_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)

    assert "Hola Mundo" in texts
    assert "Texto de Prueba" in texts
