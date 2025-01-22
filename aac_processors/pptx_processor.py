import os
import signal
from functools import wraps
from io import BytesIO
from typing import Any, Callable, Optional

import requests
from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.action import PP_ACTION  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore

from .base_processor import AACProcessor
from .tree_structure import AACButton, AACPage, AACTree, ButtonType


def timeout(seconds: int) -> Callable:
    """Create a timeout decorator."""

    def decorator(func: Callable) -> Callable:
        @wraps(func)
        def wrapper(*args: Any, **kwargs: Any) -> Any:
            def handler(signum: int, frame: Any) -> None:
                raise TimeoutError(f"Function timed out after {seconds} seconds")

            # Set the timeout handler
            old_handler = signal.signal(signal.SIGALRM, handler)
            signal.alarm(seconds)

            try:
                result = func(*args, **kwargs)
            finally:
                # Restore the old handler and disable alarm
                signal.alarm(0)
                signal.signal(signal.SIGALRM, old_handler)
            return result

        return wrapper

    return decorator


class PowerPointProcessor(AACProcessor):
    """Processor for PowerPoint files."""

    def __init__(self) -> None:
        """Initialize PowerPoint processor."""
        super().__init__()
        self.collected_texts: list[str] = []

    def can_process(self, file_path: str) -> bool:
        """Check if file can be processed."""
        return file_path.lower().endswith(".pptx")

    @timeout(30)  # 30 second timeout for image processing
    def _process_image(
        self,
        slide: Any,
        img_url: str,
        left: int,
        top: int,
        box_width: int,
        box_height: int,
    ) -> bool:
        """Process and add an image to the slide with timeout."""
        response = requests.get(img_url, timeout=10)
        if response.status_code == 200:
            img_data = BytesIO(response.content)
            # Add image to top portion of button
            img_height = box_height * 0.7
            slide.shapes.add_picture(
                img_data,
                left + box_width * 0.1,
                top + box_height * 0.1,
                width=box_width * 0.8,
                height=img_height,
            )
            return True
        return False

    def load_into_tree(self, file_path: str) -> AACTree:
        """Load PowerPoint file into tree structure."""
        tree = AACTree()
        prs = Presentation(file_path)

        # First pass: create pages and map slide indices to page IDs
        slide_to_page: dict[int, str] = {}
        for idx, slide in enumerate(prs.slides, 1):
            page_id = f"slide_{idx}"
            if idx == 1:
                tree.root_id = page_id

            title_text = (
                slide.shapes.title.text if slide.shapes.title else f"Slide {idx}"
            )
            page = AACPage(
                id=page_id, name=title_text, grid_size=(4, 3)  # Default grid size
            )
            tree.pages[page_id] = page
            slide_to_page[idx - 1] = page_id

        # Second pass: process shapes and detect navigation links
        for idx, slide in enumerate(prs.slides):
            page = tree.pages[slide_to_page[idx]]

            # Process shapes on slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and shape != slide.shapes.title:  # Skip title shape
                        # Check if shape has a hyperlink to another slide
                        button_type = ButtonType.SPEAK
                        target_page_id = None

                        if hasattr(shape, "click_action"):
                            action = shape.click_action
                            has_target = (
                                action.action == PP_ACTION.HYPERLINK
                                and hasattr(action, "target_slide")
                            )
                            if has_target:
                                target_slide = action.target_slide
                                target_idx = prs.slides.index(target_slide)
                                if target_idx in slide_to_page:
                                    button_type = ButtonType.NAVIGATE
                                    target_page_id = slide_to_page[target_idx]

                        # Create button
                        button = AACButton(
                            id=f"{page.id}_btn_{len(page.buttons) + 1}",
                            label=text,
                            vocalization=text,
                            type=button_type,
                            position=(len(page.buttons) // 3, len(page.buttons) % 3),
                        )

                        if target_page_id:
                            button.target_page_id = target_page_id

                        # Try to extract background color
                        if hasattr(shape, "fill"):
                            fill = shape.fill
                            if fill.type == 1:  # MSO_FILL.SOLID
                                if hasattr(fill.fore_color, "rgb"):
                                    rgb = fill.fore_color.rgb
                                    if rgb:
                                        # Convert RGB tuple to hex color
                                        color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                        button.style.body_color = color

                        page.buttons.append(button)

        return tree

    def save_from_tree(self, tree: AACTree, output_path: str) -> None:
        """Save tree structure to PowerPoint file."""
        try:
            # Create new presentation
            prs = Presentation()

            # Set slide size to a standard 16:9 aspect ratio
            prs.slide_width = 9144000  # 9144000 EMUs = 10 inches
            prs.slide_height = 5143500  # 5143500 EMUs = 5.625 inches

            # Create a mapping of page IDs to slide numbers for navigation
            page_to_slide: dict[str, int] = {}

            # Process each page
            for slide_idx, (page_id, page) in enumerate(tree.pages.items(), 1):
                try:
                    # Add a slide with blank layout
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    page_to_slide[page_id] = slide_idx - 1

                    # Add title to slide
                    title_box = slide.shapes.add_textbox(
                        int(prs.slide_width * 0.1),
                        int(prs.slide_height * 0.02),
                        int(prs.slide_width * 0.8),
                        int(prs.slide_height * 0.08),
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = page.name
                    for paragraph in title_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER

                    # Calculate grid cell size
                    content_left = int(prs.slide_width * 0.1)
                    # Move down to account for title
                    content_top = int(prs.slide_height * 0.15)
                    content_width = int(prs.slide_width * 0.8)
                    content_height = int(prs.slide_height * 0.75)

                    cell_width = content_width / page.grid_size[1]
                    cell_height = content_height / page.grid_size[0]

                    # Add buttons as shapes
                    for button in page.buttons:
                        try:
                            # Calculate position and size
                            left = content_left + int(button.position[1] * cell_width)
                            top = content_top + int(button.position[0] * cell_height)
                            box_width = int(cell_width * 0.9)
                            box_height = int(cell_height * 0.9)

                            # Create shape for button
                            shape = slide.shapes.add_shape(
                                1,  # MSO_SHAPE.RECTANGLE
                                left,
                                top,
                                box_width,
                                box_height,
                            )

                            # Set background color
                            shape.fill.solid()
                            if button.style.body_color:
                                try:
                                    color = button.style.body_color.lstrip("#")
                                    r = int(color[0:2], 16)
                                    g = int(color[2:4], 16)
                                    b = int(color[4:6], 16)
                                    # RGBColor constructor takes RGB values directly
                                    shape.fill.fore_color.rgb = RGBColor(r, g, b)
                                except (ValueError, IndexError):
                                    shape.fill.fore_color.rgb = RGBColor(200, 220, 255)
                            else:
                                shape.fill.fore_color.rgb = RGBColor(200, 220, 255)

                            # Add image if present
                            image_height = box_height * 0.6  # Reduce image height
                            if button.image and "url" in button.image:
                                try:
                                    img_url = button.image["url"]
                                    success = self._process_image(
                                        slide,
                                        img_url,
                                        left,
                                        top,
                                        box_width,
                                        image_height,
                                    )
                                    if not success:
                                        self.debug(f"Failed to add image: {img_url}")
                                except TimeoutError:
                                    self.debug("Image processing timed out")
                                except Exception as e:
                                    self.debug(f"Error adding image: {str(e)}")

                            # Add text to shape
                            tf = shape.text_frame
                            tf.text = button.label or ""
                            # Add margin below image
                            tf.margin_top = int(image_height + box_height * 0.05)
                            tf.margin_bottom = int(box_height * 0.05)
                            tf.margin_left = int(box_width * 0.05)
                            tf.margin_right = int(box_width * 0.05)

                            # Center align text
                            for paragraph in tf.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER

                            # Add hyperlink for navigation buttons
                            is_nav = (
                                button.type == ButtonType.NAVIGATE
                                and button.target_page_id
                            )
                            if is_nav and button.target_page_id in page_to_slide:
                                target_slide = prs.slides[
                                    page_to_slide[button.target_page_id]
                                ]
                                shape.click_action.target_slide = target_slide
                                shape.click_action.action = PP_ACTION.HYPERLINK

                                # Make navigation buttons visually distinct
                                shape.line.color.rgb = RGBColor(0, 0, 255)  # Blue
                                shape.line.width = 50800  # 4 points

                                # Add small arrow indicator
                                arrow = slide.shapes.add_shape(
                                    77,  # Right arrow shape
                                    left + box_width - int(box_width * 0.2),
                                    top + box_height - int(box_height * 0.2),
                                    int(box_width * 0.15),
                                    int(box_height * 0.15),
                                )
                                arrow.fill.solid()
                                arrow.fill.fore_color.rgb = RGBColor(0, 0, 255)

                        except Exception as e:
                            self.debug(f"Error processing button: {str(e)}")
                            continue

                except Exception as e:
                    self.debug(f"Error processing page: {str(e)}")
                    continue

            # Save the presentation
            prs.save(output_path)

        except Exception as e:
            self.debug(f"Error in save_from_tree: {str(e)}")
            raise

    def process_files(
        self, output_dir: str, translations: Optional[dict[str, str]] = None
    ) -> Optional[str]:
        """Process PowerPoint files for text extraction or translation."""
        if not self.source_file:
            return None

        source_path = str(self.source_file)

        if translations is None:
            # Extract mode
            self.extract_texts(source_path)
            return None

        # Translation mode
        if "target_lang" not in translations:
            return None

        # Create output path
        base_name = os.path.basename(source_path)
        name, ext = os.path.splitext(base_name)
        output_path = os.path.join(
            output_dir, f"{name}_{translations['target_lang']}{ext}"
        )

        # Load and translate
        tree = self.load_into_tree(source_path)

        # Translate button labels
        for page in tree.pages.values():
            for button in page.buttons:
                if button.label in translations:
                    button.label = translations[button.label]
                    button.vocalization = button.label

        # Save translated file
        self.save_from_tree(tree, output_path)

        return output_path

    def extract_texts(self, file_path: str) -> list[str]:
        """Extract all texts from PowerPoint file."""
        texts = []
        prs = Presentation(file_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)
                        self.collected_texts.append(text)

        return texts

    def create_translated_file(
        self, file_path: str, translations: dict[str, str]
    ) -> Optional[str]:
        """Create translated version of PowerPoint file.

        Args:
            file_path: Path to PowerPoint file
            translations: Dictionary of translations

        Returns:
            Path to translated file if successful
        """
        try:
            # Load file into tree
            tree = self.load_into_tree(file_path)

            # Apply translations
            for page in tree.pages.values():
                for button in page.buttons:
                    if button.label in translations:
                        button.label = translations[button.label]
                        button.vocalization = translations[button.label]

            # Create output path with language code
            target_lang = translations.get("target_lang", "translated")
            base_name = os.path.splitext(file_path)[0]
            output_path = f"{base_name}_{target_lang}.pptx"

            # Save translated tree
            self.save_from_tree(tree, output_path)
            return output_path

        except Exception as e:
            self.debug(f"Error creating translated file: {str(e)}")
            return None
